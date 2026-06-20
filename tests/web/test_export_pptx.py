import pytest
from pptx import Presentation
import io

SAMPLE_DATA = {
    "ticker": "AAPL",
    "date": "2026-06-20",
    "signal": "BUY",
    "reports": {
        "market_report": (
            "AAPL (Apple Inc.) — Technical Analysis\n"
            "## Executive Summary\n"
            "Apple is in a strong uptrend. Price above all major MAs.\n"
            "- Close (Jun 20): $200.00\n- 10 EMA: $195.00\n- 50 SMA: $185.00\n- 200 SMA: $160.00\n"
            "## Trend Analysis\n- Bullish above all MAs\n- RSI at 62 (not overbought)\n"
        ),
        "sentiment_report": (
            "**Overall Sentiment:** Bullish (Score: 7.5/10)\n"
            "**Confidence:** High\n- Strong retail interest\n- Institutional accumulation\n"
        ),
        "news_report": (
            "Top news:\n- Apple AI features driving upgrade cycle\n"
            "- Services revenue record high\n- China sales recovering\n"
        ),
        "fundamentals_report": (
            "P/E: 28x  Revenue: $400B  Net Margin: 26%\n"
            "EPS growth: 12% YoY  Debt/Equity: 1.8\n"
        ),
        "investment_plan": "Bull case: AI supercycle\nBear case: valuation stretched\n",
        "trader_investment_plan": "BUY 100 shares at $200, stop $185, target $225\n",
    },
    "debate": {
        "investment": {
            "Bull Researcher": "Strong AI momentum",
            "Bear Researcher": "Valuation risk",
            "Research Manager": "Lean bullish",
        },
        "risk": {
            "Aggressive Analyst": "Full position",
            "Conservative Analyst": "Half position",
            "Portfolio Manager": "75% position, BUY",
        },
    },
    "decision_text": (
        "BUY Apple at $200 with support at $185 and resistance at $225. "
        "AI-driven growth justifies premium valuation. Risk: macro headwinds."
    ),
}


def test_generate_pptx_returns_bytes():
    from web.export_pptx import generate_pptx
    result = generate_pptx(SAMPLE_DATA)
    assert isinstance(result, bytes)
    assert len(result) > 1000


def test_generate_pptx_valid_pptx():
    from web.export_pptx import generate_pptx
    result = generate_pptx(SAMPLE_DATA)
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 9


def test_generate_pptx_sell_signal():
    from web.export_pptx import generate_pptx
    data = {**SAMPLE_DATA, "signal": "SELL"}
    result = generate_pptx(data)
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 9
