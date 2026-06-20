"""PPT generation for trading analysis reports.

Generates a 9-slide PPTX deck from a dict of markdown report sections.
Slide order: Cover, Executive Summary, Fundamentals, News,
             Sentiment, Technical, Research Debate, Risk, Final Decision.
"""
import io
import re

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── palette ───────────────────────────────────────────────────────────────
_BG      = RGBColor(0x0F, 0x17, 0x2A)   # page background
_SURFACE = RGBColor(0x1E, 0x29, 0x3B)   # card / header background
_DEEP    = RGBColor(0x19, 0x24, 0x38)   # alternating row
_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
_MUTED   = RGBColor(0x94, 0xA3, 0xB8)
_BUY     = RGBColor(0x10, 0xB9, 0x81)
_SELL    = RGBColor(0xEF, 0x44, 0x44)
_HOLD    = RGBColor(0xF5, 0x9E, 0x0B)

_SLIDE_W     = Inches(13.333)
_SLIDE_H     = Inches(7.5)
_CONTENT_TOP = Inches(0.78)  # below the chrome header


def _signal_color(signal: str) -> RGBColor:
    s = signal.upper()
    if s in ("BUY", "OVERWEIGHT"):   return _BUY
    if s in ("SELL", "UNDERWEIGHT"): return _SELL
    return _HOLD


def _hex_signal(signal: str) -> str:
    s = signal.upper()
    if s in ("BUY", "OVERWEIGHT"):   return "#10B981"
    if s in ("SELL", "UNDERWEIGHT"): return "#EF4444"
    return "#F59E0B"


def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _BG


def _add_text(slide, text, left, top, width, height,
              size=18, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or _WHITE


def _add_chrome(slide, title: str, ticker: str, date: str, page: int,
                accent: RGBColor = None):
    """Colored header band + title + footer on every content slide."""
    col = accent or _SURFACE

    # Full-width header rectangle
    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), _SLIDE_W, Inches(0.6))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = col
    hdr.line.color.rgb = col

    tf = hdr.text_frame
    tf.margin_left = Pt(14)
    tf.margin_top  = Pt(8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(20)
    run.font.bold  = True
    run.font.color.rgb = _WHITE

    # Thin accent separator below header (slightly lighter shade)
    sep = slide.shapes.add_shape(1, Inches(0), Inches(0.6), _SLIDE_W, Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(
        min(col[0] + 30, 255),
        min(col[1] + 30, 255),
        min(col[2] + 30, 255),
    )
    sep.line.color.rgb = sep.fill.fore_color.rgb

    # Footer line
    ftr = slide.shapes.add_shape(1, Inches(0), Inches(7.26), _SLIDE_W, Inches(0.02))
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = _SURFACE
    ftr.line.color.rgb = _SURFACE

    _add_text(slide,
              f"{ticker}  ·  {date}  ·  AI Multi-Agent Analysis  ·  {page} / 9",
              Inches(0.4), Inches(7.3), Inches(12.5), Inches(0.25),
              size=7.5, color=_MUTED)


def _add_bullets(slide, items: list[str], left, top, width, height,
                 size=14, color=None, bullet_char="▸"):
    if not items:
        return
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color or _WHITE


def _extract_bullets(text: str, max_bullets: int = 5) -> list[str]:
    """Return up to max_bullets plain strings from markdown text."""
    bullets = []
    for raw in text.split('\n'):
        line = raw.strip()
        if not line or line.startswith('#'):
            continue
        line = re.sub(r'\*{1,3}([^*]+)\*{1,3}', r'\1', line)
        line = re.sub(r'`([^`]+)`', r'\1', line)
        line = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', line)
        if line.startswith(('- ', '* ', '• ', '– ')):
            bullets.append(line[2:].strip())
        elif re.match(r'^\d+[.)]\s+', line):
            bullets.append(re.sub(r'^\d+[.)]\s+', '', line))
        elif len(bullets) == 0 and len(line) > 30:
            sentences = re.split(r'(?<=[.!?])\s+', line)
            bullets.extend(s for s in sentences[:2] if len(s) > 15)
        if len(bullets) >= max_bullets:
            break
    return bullets[:max_bullets]


def _extract_score(text: str) -> float | None:
    m = re.search(r'(\d+(?:\.\d+)?)\s*/\s*10', text)
    return float(m.group(1)) if m else None


def _extract_kv_table(text: str) -> list[tuple[str, str]]:
    rows = []
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        m = re.match(r'^([A-Za-z/\s]+?)\s*[:\t]\s*(.+)$', line)
        if m:
            key = m.group(1).strip()[:30]
            val = re.sub(r'\*{1,3}([^*]+)\*{1,3}', r'\1', m.group(2).strip())[:40]
            rows.append((key, val))
        if len(rows) >= 8:
            break
    return rows


def _extract_price_targets(text: str) -> dict:
    """Extract entry, target price, and stop loss from decision text."""
    result = {}
    patterns = {
        'entry':  r'entry\s*(?:point|price|level)?[:\s]+(?:HKD|USD|[\$€¥£])?\s*([\d,]+(?:\.\d+)?)',
        'target': r'(?:price\s*)?target\s*(?:price)?[:\s]+(?:HKD|USD|[\$€¥£])?\s*([\d,]+(?:\.\d+)?)',
        'stop':   r'stop\s*(?:loss)?[:\s]+(?:HKD|USD|[\$€¥£])?\s*([\d,]+(?:\.\d+)?)',
    }
    for key, pat in patterns.items():
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            try:
                result[key] = float(m.group(1).replace(',', ''))
            except ValueError:
                pass
    return result


def _extract_company_name(ticker: str, market_report: str) -> str:
    if not market_report:
        return ticker
    escaped = re.escape(ticker)
    m = re.search(rf'{escaped}\s*\(([^)]+)\)', market_report)
    return m.group(1).strip() if m else ticker


# ── chart helpers ─────────────────────────────────────────────────────────

def _make_sentiment_gauge(score: float) -> bytes:
    fig, ax = plt.subplots(figsize=(4, 2.4), facecolor='#0F172A')
    ax.set_facecolor('#0F172A')
    ax.set_xlim(-1.3, 1.3)
    ax.set_ylim(-0.35, 1.3)
    ax.set_aspect('equal')
    ax.axis('off')

    theta = np.linspace(np.pi, 0, 200)
    ax.plot(np.cos(theta), np.sin(theta), color='#1E293B', linewidth=14)

    frac = max(0.0, min(score / 10.0, 1.0))
    theta_fill = np.linspace(np.pi, np.pi - frac * np.pi, 200)
    gauge_color = '#EF4444' if score < 4 else ('#F59E0B' if score < 7 else '#10B981')
    ax.plot(np.cos(theta_fill), np.sin(theta_fill), color=gauge_color, linewidth=14)

    angle = np.pi - frac * np.pi
    ax.annotate('', xy=(0.72 * np.cos(angle), 0.72 * np.sin(angle)),
                xytext=(0, 0),
                arrowprops=dict(arrowstyle='->', color='white', lw=2.5))

    ax.text(0, -0.28, f'{score:.1f} / 10', ha='center', va='center',
            fontsize=15, color='white', fontweight='bold')

    label = 'BEARISH' if score < 4 else ('NEUTRAL' if score < 7 else 'BULLISH')
    ax.text(0, -0.06, label, ha='center', va='center',
            fontsize=9, color=gauge_color)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=130, facecolor='#0F172A')
    plt.close(fig)
    return buf.getvalue()


def _make_risk_bar(text: str) -> bytes:
    t = text.lower()
    if 'high risk' in t or 'elevated risk' in t:
        level = 'HIGH'
    elif 'medium risk' in t or 'moderate risk' in t:
        level = 'MEDIUM'
    else:
        level = 'LOW'

    fig, ax = plt.subplots(figsize=(5, 1.6), facecolor='#0F172A')
    ax.set_facecolor('#0F172A')
    ax.axis('off')

    labels = ['LOW', 'MEDIUM', 'HIGH']
    colors = ['#10B981', '#F59E0B', '#EF4444']
    for i, (lbl, c) in enumerate(zip(labels, colors)):
        alpha = 1.0 if lbl == level else 0.22
        ax.barh(0, 1, left=i * 1.25, height=0.65, color=c, alpha=alpha)
        ax.text(i * 1.25 + 0.5, 0, lbl, ha='center', va='center',
                color='white', fontsize=11, fontweight='bold')

    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=130, facecolor='#0F172A')
    plt.close(fig)
    return buf.getvalue()


def _make_price_bar(text: str) -> bytes:
    patterns = {
        'Close':   r'close[^:\n]{0,10}[:=]\s*\$?([\d,.]+)',
        '10 EMA':  r'10\s*ema[^:\n]{0,10}[:=]\s*\$?([\d,.]+)',
        '50 SMA':  r'50\s*sma[^:\n]{0,10}[:=]\s*\$?([\d,.]+)',
        '200 SMA': r'200\s*sma[^:\n]{0,10}[:=]\s*\$?([\d,.]+)',
    }
    vals = {}
    for label, pat in patterns.items():
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            try:
                vals[label] = float(m.group(1).replace(',', ''))
            except ValueError:
                pass
    if not vals:
        return None

    fig, ax = plt.subplots(figsize=(5.2, 2.8), facecolor='#0F172A')
    ax.set_facecolor('#1E293B')
    labels = list(vals.keys())
    values = list(vals.values())
    close  = vals.get('Close', values[0])
    bar_colors = []
    for k, v in vals.items():
        if k == 'Close':
            bar_colors.append('#10B981')
        elif v <= close:
            bar_colors.append('#3B82F6')
        else:
            bar_colors.append('#64748B')

    bars = ax.barh(labels, values, color=bar_colors, height=0.55)
    for i, v in enumerate(values):
        ax.text(v * 1.001, i, f'  {v:,.1f}', va='center', color='white', fontsize=9)
    ax.set_xlabel('Price', color='#94A3B8', fontsize=9)
    ax.tick_params(colors='#94A3B8')
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.xaxis.set_tick_params(labelcolor='#94A3B8')
    ax.yaxis.set_tick_params(labelcolor='white', labelsize=9)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=130, facecolor='#0F172A')
    plt.close(fig)
    return buf.getvalue()


def _add_chart_image(slide, png_bytes: bytes, left, top, width, height):
    if not png_bytes:
        return
    slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)


# ── slide builders ────────────────────────────────────────────────────────

def _slide_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    sig    = data['signal']
    col    = _signal_color(sig)
    ticker = data['ticker']
    date   = data.get('date', '')
    company = _extract_company_name(ticker, data.get('reports', {}).get('market_report', ''))

    # Left color stripe — 33% of width
    stripe = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.4), _SLIDE_H)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = col
    stripe.line.color.rgb = col

    # Signal text on stripe
    _add_text(slide, sig,
              Inches(0.2), Inches(2.0), Inches(4.0), Inches(1.6),
              size=56, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide, "RECOMMENDATION",
              Inches(0.2), Inches(3.5), Inches(4.0), Inches(0.4),
              size=10, color=_WHITE, align=PP_ALIGN.CENTER)

    # Branding bottom-left
    _add_text(slide, "AI Multi-Agent Analysis",
              Inches(0.2), Inches(6.8), Inches(4.0), Inches(0.4),
              size=8.5, color=_WHITE, align=PP_ALIGN.CENTER)

    # Right content area — ticker
    _add_text(slide, ticker,
              Inches(4.8), Inches(1.3), Inches(8.2), Inches(1.3),
              size=56, bold=True)

    if company and company != ticker:
        _add_text(slide, company,
                  Inches(4.8), Inches(2.7), Inches(8.2), Inches(0.7),
                  size=20, color=_MUTED)

    # Thin divider
    div = slide.shapes.add_shape(1, Inches(4.8), Inches(3.65), Inches(8.2), Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = _SURFACE
    div.line.color.rgb = _SURFACE

    _add_text(slide, f"Analysis Date:  {date}",
              Inches(4.8), Inches(3.82), Inches(8.2), Inches(0.5),
              size=15, color=_MUTED)
    _add_text(slide, "Market  ·  Sentiment  ·  News  ·  Fundamentals  ·  Research  ·  Risk",
              Inches(4.8), Inches(4.5), Inches(8.2), Inches(0.45),
              size=11, color=_MUTED)

    _add_text(slide,
              "For informational purposes only. Not investment advice.",
              Inches(4.8), Inches(7.0), Inches(8.2), Inches(0.4),
              size=9, color=_MUTED)


def _slide_exec_summary(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    sig  = data.get('signal', 'HOLD')
    col  = _signal_color(sig)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "Executive Summary", tick, dt, 2, col)

    # Main bullets — prefer decision_text which has the actual analyst conclusion
    decision = data.get('decision_text', '')
    market   = data.get('reports', {}).get('market_report', '')
    src      = decision if len(decision) > 100 else market
    bullets  = _extract_bullets(src, max_bullets=5)
    _add_bullets(slide, bullets,
                 Inches(0.4), _CONTENT_TOP, Inches(8.5), Inches(6.2),
                 size=14)

    # ── Key Stats sidebar ──────────────────────────────────────────────────
    targets = _extract_price_targets(decision)
    stats = [
        ("Signal",   sig),
        ("Ticker",   tick),
        ("Date",     dt),
    ]
    if targets.get('target'): stats.append(("Price Target", f"{targets['target']:,.0f}"))
    if targets.get('entry'):  stats.append(("Entry Point",  f"{targets['entry']:,.0f}"))
    if targets.get('stop'):   stats.append(("Stop Loss",    f"{targets['stop']:,.0f}"))

    y = _CONTENT_TOP
    for label, value in stats[:6]:
        is_signal = (label == "Signal")
        box = slide.shapes.add_shape(1, Inches(9.2), y, Inches(3.9), Inches(0.72))
        box.fill.solid()
        box.fill.fore_color.rgb = _SURFACE
        box.line.color.rgb      = col if is_signal else RGBColor(0x2D, 0x3B, 0x55)
        box.line.width          = Pt(2 if is_signal else 0.75)

        tf = box.text_frame
        tf.margin_left = Pt(10)
        tf.margin_top  = Pt(5)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text       = f"{label}:  {value}"
        run.font.size  = Pt(13 if is_signal else 11)
        run.font.bold  = is_signal
        run.font.color.rgb = col if is_signal else _WHITE
        y += Inches(0.78)


def _slide_fundamentals(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "Fundamentals", tick, dt, 3)

    text = data.get('reports', {}).get('fundamentals_report', '')
    rows = _extract_kv_table(text)

    if rows:
        n_rows = len(rows)
        row_h  = min(0.46, 5.8 / (n_rows + 1))
        tbl_h  = row_h * (n_rows + 1)
        tbl = slide.shapes.add_table(
            n_rows + 1, 2,
            Inches(0.4), _CONTENT_TOP,
            Inches(5.6), Inches(tbl_h)
        ).table

        # Header row
        for ci, header in enumerate(('Metric', 'Value')):
            cell = tbl.cell(0, ci)
            cell.text = header
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.bold  = True
            p.runs[0].font.size  = Pt(11)
            p.runs[0].font.color.rgb = _WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = _SURFACE

        # Data rows with alternating background
        for ri, (k, v) in enumerate(rows):
            bg = _BG if ri % 2 == 0 else _DEEP
            for ci, val in enumerate((k, v)):
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.runs[0].font.size      = Pt(10)
                p.runs[0].font.color.rgb = _WHITE if ci == 1 else _MUTED
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
    else:
        bullets = _extract_bullets(text, max_bullets=7)
        _add_bullets(slide, bullets,
                     Inches(0.4), _CONTENT_TOP, Inches(12.5), Inches(6.2))


def _slide_news(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "News Analysis", tick, dt, 4)

    text    = data.get('reports', {}).get('news_report', '')
    bullets = _extract_bullets(text, max_bullets=6)
    _add_bullets(slide, bullets,
                 Inches(0.4), _CONTENT_TOP, Inches(12.5), Inches(6.2),
                 size=14)


def _slide_sentiment(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "Sentiment Analysis", tick, dt, 5)

    text    = data.get('reports', {}).get('sentiment_report', '')
    bullets = _extract_bullets(text, max_bullets=5)
    _add_bullets(slide, bullets,
                 Inches(0.4), _CONTENT_TOP, Inches(7.9), Inches(6.2),
                 size=14)

    score = _extract_score(text)
    if score is not None:
        png = _make_sentiment_gauge(score)
        _add_chart_image(slide, png, Inches(8.6), _CONTENT_TOP, Inches(4.4), Inches(2.8))


def _slide_technical(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "Technical Analysis", tick, dt, 6)

    text    = data.get('reports', {}).get('market_report', '')
    bullets = _extract_bullets(text, max_bullets=5)
    _add_bullets(slide, bullets,
                 Inches(0.4), _CONTENT_TOP, Inches(7.9), Inches(4.6),
                 size=14)

    png = _make_price_bar(text)
    if png:
        _add_chart_image(slide, png, Inches(8.3), _CONTENT_TOP, Inches(4.7), Inches(3.2))


def _slide_research_debate(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "Research Debate", tick, dt, 7)

    invest = data.get('debate', {}).get('investment', {})

    # ── Bull column ────────────────────────────────────────────────────────
    bull_hdr = slide.shapes.add_shape(1, Inches(0.4), _CONTENT_TOP,
                                      Inches(5.9), Inches(0.42))
    bull_hdr.fill.solid()
    bull_hdr.fill.fore_color.rgb = RGBColor(0x06, 0x4E, 0x3B)
    bull_hdr.line.color.rgb      = _BUY
    bull_hdr.line.width          = Pt(1)
    tf = bull_hdr.text_frame
    tf.margin_left = Pt(10)
    tf.margin_top  = Pt(5)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "▲  Bull Case"
    run.font.size  = Pt(13)
    run.font.bold  = True
    run.font.color.rgb = _BUY

    bull_text = invest.get('Bull Researcher', '')
    _add_bullets(slide, _extract_bullets(bull_text, 4),
                 Inches(0.4), Inches(1.35), Inches(5.9), Inches(4.8),
                 size=12)

    # Vertical divider
    vdiv = slide.shapes.add_shape(1, Inches(6.5), _CONTENT_TOP,
                                  Inches(0.04), Inches(5.8))
    vdiv.fill.solid()
    vdiv.fill.fore_color.rgb = _SURFACE
    vdiv.line.color.rgb      = _SURFACE

    # ── Bear column ────────────────────────────────────────────────────────
    bear_hdr = slide.shapes.add_shape(1, Inches(6.7), _CONTENT_TOP,
                                      Inches(6.3), Inches(0.42))
    bear_hdr.fill.solid()
    bear_hdr.fill.fore_color.rgb = RGBColor(0x4C, 0x0E, 0x0E)
    bear_hdr.line.color.rgb      = _SELL
    bear_hdr.line.width          = Pt(1)
    tf2 = bear_hdr.text_frame
    tf2.margin_left = Pt(10)
    tf2.margin_top  = Pt(5)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = "▼  Bear Case"
    run2.font.size  = Pt(13)
    run2.font.bold  = True
    run2.font.color.rgb = _SELL

    bear_text = invest.get('Bear Researcher', '')
    _add_bullets(slide, _extract_bullets(bear_text, 4),
                 Inches(6.7), Inches(1.35), Inches(6.3), Inches(4.8),
                 size=12)

    # Research Manager verdict bar at bottom
    judge = invest.get('Research Manager', '')
    if judge:
        judge_bullets = _extract_bullets(judge, 1)
        if judge_bullets:
            jbar = slide.shapes.add_shape(1, Inches(0.4), Inches(6.45),
                                          Inches(12.6), Inches(0.55))
            jbar.fill.solid()
            jbar.fill.fore_color.rgb = _SURFACE
            jbar.line.color.rgb      = _SURFACE
            tf3 = jbar.text_frame
            tf3.margin_left = Pt(12)
            tf3.margin_top  = Pt(7)
            p3 = tf3.paragraphs[0]
            run3 = p3.add_run()
            run3.text = f"Research Manager:  {judge_bullets[0]}"
            run3.font.size      = Pt(11)
            run3.font.color.rgb = _MUTED


def _slide_risk(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "Risk Assessment", tick, dt, 8, _SELL)

    risk    = data.get('debate', {}).get('risk', {})
    pm_text = risk.get('Portfolio Manager', '')
    src     = pm_text or data.get('decision_text', '')
    bullets = _extract_bullets(src, max_bullets=5)
    _add_bullets(slide, bullets,
                 Inches(0.4), _CONTENT_TOP, Inches(7.9), Inches(5.5),
                 size=14)

    dec_text = data.get('decision_text', '')
    png = _make_risk_bar(dec_text + ' ' + pm_text)
    if png:
        _add_chart_image(slide, png, Inches(8.4), _CONTENT_TOP, Inches(4.6), Inches(1.9))


def _slide_final(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    sig  = data['signal']
    col  = _signal_color(sig)
    tick = data.get('ticker', '')
    dt   = data.get('date', '')
    _add_chrome(slide, "Final Decision", tick, dt, 9, col)

    dec_text = data.get('decision_text', '')
    targets  = _extract_price_targets(dec_text)

    # ── Signal badge + Ticker ──────────────────────────────────────────────
    badge = slide.shapes.add_shape(1, Inches(0.4), _CONTENT_TOP,
                                   Inches(3.6), Inches(1.1))
    badge.fill.solid()
    badge.fill.fore_color.rgb = col
    badge.line.color.rgb      = col
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = sig
    run.font.size  = Pt(38)
    run.font.bold  = True
    run.font.color.rgb = _WHITE

    _add_text(slide, tick,
              Inches(4.3), _CONTENT_TOP, Inches(9.0), Inches(0.85),
              size=30, bold=True)

    # ── Price target callout boxes ─────────────────────────────────────────
    target_items = []
    if targets.get('entry'):  target_items.append(("Entry Point",  f"{targets['entry']:,.0f}",  _MUTED))
    if targets.get('target'): target_items.append(("Price Target", f"{targets['target']:,.0f}", col))
    if targets.get('stop'):   target_items.append(("Stop Loss",    f"{targets['stop']:,.0f}",   _SELL))

    if target_items:
        x_positions = [Inches(0.4), Inches(4.6), Inches(8.8)]
        box_w = Inches(3.85)
        for i, (label, value, accent) in enumerate(target_items[:3]):
            x = x_positions[i]
            b = slide.shapes.add_shape(1, x, Inches(2.1), box_w, Inches(0.95))
            b.fill.solid()
            b.fill.fore_color.rgb = _SURFACE
            b.line.color.rgb      = accent
            b.line.width          = Pt(2)
            _add_text(slide, label,
                      x + Inches(0.15), Inches(2.14),
                      box_w - Inches(0.3), Inches(0.28),
                      size=9, color=_MUTED)
            _add_text(slide, value,
                      x + Inches(0.15), Inches(2.38),
                      box_w - Inches(0.3), Inches(0.58),
                      size=22, bold=True, color=accent)
        y_bullets = Inches(3.25)
    else:
        y_bullets = Inches(2.1)

    bullets = _extract_bullets(dec_text, max_bullets=5)
    _add_bullets(slide, bullets,
                 Inches(0.4), y_bullets, Inches(12.6), Inches(3.9),
                 size=13)

    _add_text(slide,
              "This report is for informational purposes only and does not constitute investment advice.",
              Inches(0.4), Inches(7.0), Inches(12.6), Inches(0.45),
              size=8, color=_MUTED)


# ── public API ────────────────────────────────────────────────────────────

def generate_pptx(data: dict) -> bytes:
    """Build a 9-slide PPTX and return the raw bytes."""
    prs = Presentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H

    _slide_cover(prs, data)
    _slide_exec_summary(prs, data)
    _slide_fundamentals(prs, data)
    _slide_news(prs, data)
    _slide_sentiment(prs, data)
    _slide_technical(prs, data)
    _slide_research_debate(prs, data)
    _slide_risk(prs, data)
    _slide_final(prs, data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
